"""PDFShell conversion service — self-hosted, all open-source.

Honest engine notes:
  * LibreOffice converts Office → PDF (and Office → Office) with high fidelity,
    but CANNOT convert PDF → Office (it imports PDFs into Draw, which can't
    export Writer/Calc/Impress formats). So PDF → Office uses purpose-built libs:
      - PDF → Word  : pdf2docx (text + layout reconstruction)
      - PDF → PPTX  : PyMuPDF renders each page to an image, one per slide
      - PDF → XLSX  : pdfplumber extracts tables → openpyxl
      - PDF → TXT   : pdfplumber text extraction

Contract (called by apps/web/lib/libreoffice.ts):
  GET  /health               → 200 "ok"
  POST /convert?target=<ext> → converted file bytes
       headers: x-source-ext: pdf | docx | xlsx | pptx | ...
       body:    raw source file bytes
"""

import base64
import json
import os
import re
import shutil
import subprocess
import tempfile
import threading
import time
import urllib.request
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from urllib.parse import urlparse, parse_qs

PORT = int(os.environ.get("PORT", "3001"))

# Hardening: cap request size (a body is held in memory for the duration of a
# job) and the number of concurrent heavy jobs (LibreOffice/Ghostscript each
# take hundreds of MB — unbounded parallelism OOMs small hosts).
MAX_BODY_BYTES = int(os.environ.get("PDFSHELL_MAX_BODY_MB", "200")) * 1024 * 1024
_JOBS = threading.Semaphore(int(os.environ.get("PDFSHELL_MAX_JOBS", "2")))

# Transient working directory. User files live here ONLY for the duration of a
# job — each request uses an auto-deleted temp subdir, and a janitor sweeps any
# stragglers (e.g. after a crash). Nothing is ever persisted.
TMP_ROOT = os.environ.get("PDFSHELL_TMP", os.path.join(tempfile.gettempdir(), "pdfshell-work"))
TMP_TTL_SECS = int(os.environ.get("PDFSHELL_TMP_TTL", "1800"))  # 30 min safety net
os.makedirs(TMP_ROOT, exist_ok=True)


def _janitor() -> None:
    """Periodically delete anything in TMP_ROOT older than the TTL."""
    while True:
        try:
            cutoff = time.time() - TMP_TTL_SECS
            for name in os.listdir(TMP_ROOT):
                path = os.path.join(TMP_ROOT, name)
                try:
                    if os.path.getmtime(path) < cutoff:
                        if os.path.isdir(path):
                            shutil.rmtree(path, ignore_errors=True)
                        else:
                            os.remove(path)
                except OSError:
                    pass
        except OSError:
            pass
        time.sleep(300)

# Self-hosted local LLM (Ollama) for opt-in AI enhancement. Nothing leaves the
# host: the convert service talks to Ollama over the local network.
OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://ollama:11434")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "llama3.2:1b")
# Capped context keeps memory low; num_gpu=0 forces CPU (safe default for hosts
# with little/no spare VRAM). On a GPU host, set OLLAMA_NUM_GPU=-1 to use it.
OLLAMA_NUM_CTX = int(os.environ.get("OLLAMA_NUM_CTX", "4096"))
OLLAMA_NUM_GPU = int(os.environ.get("OLLAMA_NUM_GPU", "0"))
# Bound the AI work so a slow CPU model can't run long enough for the request to
# time out — if it exceeds this, we fall back to faithful image slides.
OLLAMA_NUM_PREDICT = int(os.environ.get("OLLAMA_NUM_PREDICT", "1024"))
OLLAMA_TIMEOUT = int(os.environ.get("OLLAMA_TIMEOUT", "90"))

# The app exposes the admin-selected active model here (internal, token-gated).
# If reachable, it overrides the env Ollama defaults so the WHOLE app uses
# whatever the admin activated in the AI Models page.
APP_URL = os.environ.get("APP_URL", "http://host.docker.internal:4321")
INTERNAL_TOKEN = os.environ.get("PDFSHELL_INTERNAL_TOKEN", "pdfshell-internal")

MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "txt": "text/plain",
}

CORS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "POST, GET, OPTIONS",
    "Access-Control-Allow-Headers": "content-type, x-source-ext, x-password, x-owner-password, x-permissions",
}


def office_to_pdf_or_office(in_path: str, out_dir: str, target: str) -> str:
    """LibreOffice headless conversion. Returns the produced file path."""
    subprocess.run(
        [
            "soffice",
            "--headless",
            "--norestore",
            f"-env:UserInstallation=file://{os.path.join(out_dir, 'profile')}",
            "--convert-to",
            target,
            "--outdir",
            out_dir,
            in_path,
        ],
        check=True,
        timeout=180,
        capture_output=True,
    )
    base = os.path.splitext(os.path.basename(in_path))[0]
    produced = os.path.join(out_dir, f"{base}.{target}")
    if not os.path.isfile(produced):
        raise RuntimeError(f"LibreOffice produced no {target} output")
    return produced


def pdf_to_docx(in_path: str, out_path: str) -> None:
    from pdf2docx import Converter

    cv = Converter(in_path)
    try:
        cv.convert(out_path)
    finally:
        cv.close()


# Standard landscape 16:9 deck dimensions (EMU).
SLIDE_W = 12192000
SLIDE_H = 6858000


def ollama_available() -> bool:
    try:
        with urllib.request.urlopen(f"{OLLAMA_URL}/api/tags", timeout=3) as r:
            return r.status == 200
    except Exception:
        return False


def _post_json(url: str, payload: dict, headers: dict, timeout: int = 300) -> dict:
    req = urllib.request.Request(
        url, data=json.dumps(payload).encode(), headers={"content-type": "application/json", **headers}
    )
    with urllib.request.urlopen(req, timeout=timeout) as r:
        return json.loads(r.read())


def fetch_active_model():
    """Return the admin-selected active model {provider, baseUrl, model, apiKey},
    or None to fall back to the env Ollama defaults."""
    try:
        req = urllib.request.Request(
            f"{APP_URL.rstrip('/')}/api/ai/active", headers={"x-internal-token": INTERNAL_TOKEN}
        )
        with urllib.request.urlopen(req, timeout=4) as r:
            return json.loads(r.read()).get("active")
    except Exception:
        return None


def llm_json(prompt: str) -> dict:
    """Run a JSON-returning completion against the active model (or env Ollama).
    Supports Ollama, OpenAI-compatible, and Anthropic. Returns {} on failure."""
    cfg = fetch_active_model() or {
        "provider": "ollama",
        "baseUrl": OLLAMA_URL,
        "model": OLLAMA_MODEL,
        "apiKey": "",
    }
    provider, base, model, key = cfg["provider"], cfg["baseUrl"].rstrip("/"), cfg["model"], cfg.get("apiKey", "")
    try:
        if provider == "ollama":
            resp = _post_json(
                f"{base}/api/generate",
                {"model": model, "prompt": prompt, "stream": False, "format": "json",
                 "keep_alive": "30m",  # keep model resident (esp. on GPU) for fast repeats
                 "options": {"temperature": 0.2, "num_ctx": OLLAMA_NUM_CTX,
                             "num_gpu": OLLAMA_NUM_GPU, "num_predict": OLLAMA_NUM_PREDICT}},
                {},
                timeout=OLLAMA_TIMEOUT,
            )
            return json.loads(resp.get("response", "{}"))
        if provider == "anthropic":
            resp = _post_json(
                f"{base}/v1/messages",
                {"model": model, "max_tokens": 2048, "messages": [{"role": "user", "content": prompt}]},
                {"x-api-key": key, "anthropic-version": "2023-06-01"},
            )
            text = "".join(b.get("text", "") for b in resp.get("content", []))
            return json.loads(text or "{}")
        # openai / custom (OpenAI-compatible chat completions)
        resp = _post_json(
            f"{base}/v1/chat/completions",
            {"model": model, "temperature": 0.2, "response_format": {"type": "json_object"},
             "messages": [{"role": "user", "content": prompt}]},
            {"authorization": f"Bearer {key}"},
        )
        return json.loads(resp["choices"][0]["message"]["content"])
    except Exception:
        return {}


def ai_slides_from_text(full_text: str):
    """Turn document text into a clean presentation outline via the active model.
    Returns a list of {title, bullets[]} or [] on failure (caller falls back)."""
    prompt = (
        "You convert document text into a concise slide deck outline. "
        "Group related content into logical slides. Each bullet is short "
        "(max ~12 words); 3-8 bullets per slide.\n"
        'Respond with ONLY a JSON object of this exact shape:\n'
        '{"slides":[{"title":"string","bullets":["string","string"]}]}\n\n'
        f"DOCUMENT:\n{full_text[:4500]}"
    )
    data = llm_json(prompt)
    if not data:
        return []

    # Models vary in shape under format=json. Accept {"slides":[...]}, a bare
    # list, or the first list-valued key in the object.
    raw_slides = None
    if isinstance(data, dict):
        if isinstance(data.get("slides"), list):
            raw_slides = data["slides"]
        else:
            for v in data.values():
                if isinstance(v, list):
                    raw_slides = v
                    break
    elif isinstance(data, list):
        raw_slides = data
    if not raw_slides:
        return []

    def as_bullets(value):
        if isinstance(value, list):
            return [str(b).strip() for b in value if str(b).strip()]
        if isinstance(value, str) and value.strip():
            return [value.strip()]
        return []

    out = []
    for s in raw_slides:
        if not isinstance(s, dict):
            if isinstance(s, str) and s.strip():
                out.append({"title": s.strip()[:160], "bullets": []})
            continue
        title = str(s.get("title") or s.get("heading") or s.get("name") or "").strip()
        bullets = as_bullets(s.get("bullets") or s.get("points") or s.get("content") or [])
        if title or bullets:
            out.append({"title": title, "bullets": bullets})
    return out


def _add_ai_slide(prs, layout, title, bullets):
    from pptx.util import Emu, Pt
    from pptx.enum.text import MSO_AUTO_SIZE

    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Emu(686000), Emu(457000), Emu(SLIDE_W - 1372000), Emu(1200000))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (title or "Slide")[:160]
    p.font.size = Pt(32)
    p.font.bold = True

    if bullets:
        bb = slide.shapes.add_textbox(Emu(686000), Emu(1850000), Emu(SLIDE_W - 1372000), Emu(SLIDE_H - 2400000))
        bf = bb.text_frame
        bf.word_wrap = True
        bf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for i, b in enumerate(bullets):
            par = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            par.text = f"•  {b}"
            par.font.size = Pt(18)


def _add_image_slide(prs, layout, page, work_dir, idx):
    """Render the page faithfully as a high-res image, centred on a landscape
    slide. Preserves the document's real layout (headings, tables, cards) —
    the right default unless the user opts into AI re-authoring."""
    from pptx.util import Emu

    pix = page.get_pixmap(dpi=170)
    img_path = os.path.join(work_dir, f"p{idx}.png")
    pix.save(img_path)
    scale = min(SLIDE_W / pix.width, SLIDE_H / pix.height)
    w, h = int(pix.width * scale), int(pix.height * scale)
    slide = prs.slides.add_slide(layout)
    slide.shapes.add_picture(
        img_path, Emu((SLIDE_W - w) // 2), Emu((SLIDE_H - h) // 2), width=Emu(w), height=Emu(h)
    )


def pdf_to_pptx(in_path: str, out_path: str, work_dir: str, use_ai: bool = False) -> None:
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(in_path)
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]

    # AI path: re-author the whole document into a clean, grouped, EDITABLE deck.
    # Falls back to faithful page images if the model is unavailable / unhelpful.
    if use_ai:
        full_text = "\n".join(page.get_text("text") for page in doc).strip()
        slides = ai_slides_from_text(full_text) if full_text else []
        if slides:
            for s in slides:
                _add_ai_slide(prs, blank, s["title"], s["bullets"])
            doc.close()
            prs.save(out_path)
            return

    # Default (no AI): one faithful page-image per slide. This keeps the
    # document looking like itself — far better than dumping every text line as a
    # bullet, which mangles formatted reports, tables and metric cards.
    for i, page in enumerate(doc):
        _add_image_slide(prs, blank, page, work_dir, i)

    doc.close()
    prs.save(out_path)


_NUM_RE = re.compile(r"^[-+]?\(?\s?[$€£]?\s?\d{1,3}(?:[,\s]\d{3})*(?:\.\d+)?\)?\s?%?$")


def _typed(value: str):
    """Return (value, number_format) — coercing numeric-looking text to numbers
    so the spreadsheet stays computable, matching how the cell reads."""
    s = (value or "").strip()
    if not s or not _NUM_RE.match(s):
        return s, None
    pct = s.endswith("%")
    neg = s.startswith("(") and s.endswith(")")
    cleaned = re.sub(r"[^\d.]", "", s)
    if cleaned in ("", "."):
        return s, None
    try:
        num = float(cleaned)
    except ValueError:
        return s, None
    if neg:
        num = -num
    if pct:
        return num / 100, "0.0%"
    has_decimal = "." in cleaned
    return num, ("#,##0.00" if has_decimal else "#,##0")


def _reconstruct_grid(page):
    """Approximate the page's visual layout as a row/column grid from word
    positions — so prose, forms and borderless tables come out close to the
    document rather than as one flat column."""
    words = page.extract_words(x_tolerance=1.5, y_tolerance=3, keep_blank_chars=False)
    if not words:
        return []

    # Group words into rows by vertical position.
    words.sort(key=lambda w: (round(w["top"], 1), w["x0"]))
    rows, current, top = [], [], None
    for w in words:
        if top is None or abs(w["top"] - top) <= 5:
            current.append(w)
            top = w["top"] if top is None else top
        else:
            rows.append(current)
            current, top = [w], w["top"]
    if current:
        rows.append(current)

    # Derive column starts from clustered left edges across the whole page.
    starts = sorted(round(w["x0"]) for w in words)
    cols = []
    for x in starts:
        if not cols or x - cols[-1] > 24:  # new column when the gap is wide
            cols.append(x)

    grid = []
    for row in rows:
        cells = [""] * len(cols)
        for w in sorted(row, key=lambda w: w["x0"]):
            ci = 0
            for j, cx in enumerate(cols):
                if w["x0"] >= cx - 6:
                    ci = j
                else:
                    break
            cells[ci] = (cells[ci] + " " + w["text"]).strip()
        grid.append(cells)
    return grid


def _hex6_from_int(c) -> str:
    return "%06X" % (int(c) & 0xFFFFFF)


def _hex6_from_rgbf(t) -> str:
    r, g, b = (max(0, min(255, int(round(x * 255)))) for x in t[:3])
    return "%02X%02X%02X" % (r, g, b)


def _xlsx_fills(page):
    """Filled rectangles on the page (section bands, status badges) as
    (Rect, hex). Skips white and full-page background fills."""
    import fitz

    out = []
    area = page.rect.width * page.rect.height or 1
    for d in page.get_drawings():
        fill = d.get("fill")
        rect = d.get("rect")
        if fill is None or rect is None:
            continue
        hx = _hex6_from_rgbf(fill)
        if hx == "FFFFFF":
            continue
        r = fitz.Rect(rect)
        if r.get_area() > 0.8 * area:  # page background, not a band
            continue
        out.append((r, hx))
    return out


def _xlsx_page(ws, page, ppage) -> None:
    """Reconstruct one page into a styled sheet: a global column grid from text
    x-positions, per-cell font/colour from the spans, fills from drawn rectangles
    (section bands & badges), borders inside ruled tables, merged wide headers,
    and column widths from the page geometry."""
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    thin = Side(style="thin", color="BFBFBF")
    border_all = Border(left=thin, right=thin, top=thin, bottom=thin)

    spans_all = []
    for block in page.get_text("dict").get("blocks", []):
        if block.get("type") != 0:
            continue
        for ln in block.get("lines", []):
            for s in ln.get("spans", []):
                if s.get("text", "").strip():
                    spans_all.append(s)
    if not spans_all:
        for t in (page.get_text("text") or "").splitlines():
            ws.append([t])
        return

    # Group spans into rows by vertical position — same-baseline text from
    # different draw calls (e.g. a label and its value) belongs on one row.
    spans_all.sort(key=lambda s: ((s["bbox"][1] + s["bbox"][3]) / 2, s["bbox"][0]))
    rows, cur, cy = [], [], None
    for s in spans_all:
        mid = (s["bbox"][1] + s["bbox"][3]) / 2
        tol = max(3.0, (s["bbox"][3] - s["bbox"][1]) * 0.6)
        if cy is None or abs(mid - cy) <= tol:
            cur.append(s)
            cy = mid if cy is None else cy
        else:
            rows.append(cur)
            cur, cy = [s], mid
    if cur:
        rows.append(cur)

    # Column starts: cluster span left edges across the whole page.
    starts = []
    for x in sorted(s["bbox"][0] for s in spans_all):
        if not starts or x - starts[-1] > 10:
            starts.append(x)
    ncols = len(starts)
    pagew = page.rect.width

    def col_of(x):
        ci = 0
        for j, sx in enumerate(starts):
            if x >= sx - 3:
                ci = j
            else:
                break
        return ci

    fills = _xlsx_fills(page)

    def fill_at(bbox):
        cx, cy_ = (bbox[0] + bbox[2]) / 2, (bbox[1] + bbox[3]) / 2
        for r, hx in fills:
            if r.x0 <= cx <= r.x1 and r.y0 <= cy_ <= r.y1:
                return hx, (r.x1 - r.x0)
        return None, 0

    tboxes = []
    if ppage is not None:
        try:
            tboxes = [t.bbox for t in ppage.find_tables()]
        except Exception:
            tboxes = []

    def in_table(bbox):
        cx, cy_ = (bbox[0] + bbox[2]) / 2, (bbox[1] + bbox[3]) / 2
        return any(x0 <= cx <= x1 and t <= cy_ <= b for (x0, t, x1, b) in tboxes)

    maxlen = {}  # column → widest text, for width fitting
    for r, spans in enumerate(rows, start=1):
        cells = {}
        for s in sorted(spans, key=lambda s: s["bbox"][0]):
            c = col_of(s["bbox"][0])
            if c in cells:
                cells[c]["text"] += " " + s["text"].strip()
                bb, nb = cells[c]["bbox"], s["bbox"]
                cells[c]["bbox"] = [min(bb[0], nb[0]), min(bb[1], nb[1]), max(bb[2], nb[2]), max(bb[3], nb[3])]
            else:
                cells[c] = {"text": s["text"].strip(), "span": s, "bbox": list(s["bbox"])}

        for c, info in cells.items():
            s = info["span"]
            flags, name = int(s.get("flags", 0)), (s.get("font", "") or "").lower()
            val, fmt = _typed(info["text"])
            cell = ws.cell(row=r, column=c + 1, value=val)
            if fmt:
                cell.number_format = fmt
                cell.alignment = Alignment(horizontal="right")
            cell.font = Font(
                bold=bool(flags & 16) or "bold" in name or "black" in name or "semibold" in name,
                italic=bool(flags & 2) or "italic" in name or "oblique" in name,
                color=_hex6_from_int(s.get("color", 0)),
                size=max(8, min(20, round(s.get("size", 11)))),
            )
            fhex, _fw = fill_at(info["bbox"])
            if fhex:
                cell.fill = PatternFill("solid", fgColor=fhex)
            if in_table(info["bbox"]):
                cell.border = border_all
            maxlen[c] = max(maxlen.get(c, 0), len(info["text"]))

        # Merge a wide single-cell row (title / section band) across the columns.
        if len(cells) == 1 and ncols > 1:
            (c, info), = cells.items()
            fhex, fw = fill_at(info["bbox"])
            wide = (info["bbox"][2] - info["bbox"][0]) > pagew * 0.4 or fw > pagew * 0.4
            if wide:
                src = ws.cell(row=r, column=c + 1)
                v, f = src.value, src.font
                if c != 0:
                    src.value = None
                    ws.cell(row=r, column=1, value=v).font = f
                ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=ncols)
                if fhex:
                    for cc in range(1, ncols + 1):
                        ws.cell(row=r, column=cc).fill = PatternFill("solid", fgColor=fhex)

    for j in range(ncols):
        x1 = starts[j + 1] if j + 1 < ncols else pagew
        geom = (x1 - starts[j]) / 6.5
        ws.column_dimensions[get_column_letter(j + 1)].width = max(8, min(64, max(geom, maxlen.get(j, 0) + 1)))


def pdf_to_xlsx(in_path: str, out_path: str) -> None:
    """High-fidelity PDF → Excel: preserves layout, fonts/colours, section bands,
    table borders and column widths (PyMuPDF for style + drawings, pdfplumber for
    ruled-table detection). One sheet per page; falls back to flat text per page."""
    import fitz
    from openpyxl import Workbook

    try:
        import pdfplumber
        plumb = pdfplumber.open(in_path)
    except Exception:
        plumb = None

    doc = fitz.open(in_path)
    wb = Workbook()
    wb.remove(wb.active)
    try:
        for pno in range(doc.page_count):
            ws = wb.create_sheet(title=f"Page {pno + 1}"[:31])
            ppage = plumb.pages[pno] if plumb and pno < len(plumb.pages) else None
            try:
                _xlsx_page(ws, doc[pno], ppage)
            except Exception:
                for t in (doc[pno].get_text("text") or "").splitlines():
                    ws.append([t])
    finally:
        doc.close()
        if plumb:
            plumb.close()

    if not wb.sheetnames:
        wb.create_sheet(title="Empty")
    wb.save(out_path)


def ghostscript_compress(in_path: str, out_path: str, preset: str) -> None:
    """Real PDF compression via Ghostscript. Reliably downsamples images and
    rewrites the file; far stronger than a pdf-lib re-save."""
    gs_preset = {
        "screen": "/screen",     # 72 dpi  — smallest
        "ebook": "/ebook",       # 150 dpi — balanced (default)
        "printer": "/printer",   # 300 dpi — high quality
        "prepress": "/prepress",
    }.get(preset, "/ebook")
    subprocess.run(
        [
            "gs",
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.5",
            f"-dPDFSETTINGS={gs_preset}",
            "-dNOPAUSE",
            "-dQUIET",
            "-dBATCH",
            "-dDetectDuplicateImages=true",
            "-dCompressFonts=true",
            "-dSubsetFonts=true",
            f"-sOutputFile={out_path}",
            in_path,
        ],
        check=True,
        timeout=180,
        capture_output=True,
    )


def pdf_to_txt(in_path: str, out_path: str) -> None:
    import pdfplumber

    with pdfplumber.open(in_path) as pdf:
        text = "\n\n".join((p.extract_text() or "") for p in pdf.pages)
    with open(out_path, "w", encoding="utf-8") as fh:
        fh.write(text)


class WrongPassword(Exception):
    """The supplied password does not open this PDF."""


def _perm_mask(allowed) -> int:
    """Build a PDF permission bitmask from a set of allowed actions. Screen-reader
    accessibility is always permitted. Note these flags are advisory — they're
    honoured by compliant viewers, not enforced cryptographically."""
    import fitz

    mask = int(fitz.PDF_PERM_ACCESSIBILITY)
    if "print" in allowed:
        mask |= int(fitz.PDF_PERM_PRINT) | int(fitz.PDF_PERM_PRINT_HQ)
    if "copy" in allowed:
        mask |= int(fitz.PDF_PERM_COPY)
    if "modify" in allowed:
        mask |= int(fitz.PDF_PERM_MODIFY) | int(fitz.PDF_PERM_ASSEMBLE)
    if "annotate" in allowed:
        mask |= int(fitz.PDF_PERM_ANNOTATE) | int(fitz.PDF_PERM_FORM)
    return mask


_ALL_PERMS = {"print", "copy", "modify", "annotate"}


def protect_pdf(body: bytes, user_pw: str, owner_pw=None, allowed=None) -> bytes:
    """Encrypt a PDF with AES-256.

    user_pw  — required to OPEN the document.
    owner_pw — required to CHANGE permissions/print settings; defaults to the
               user password. For permission restrictions to actually bite, the
               owner password must differ from the user password (otherwise the
               opener is treated as the owner and gets full rights).
    allowed  — set of permitted actions; None means "all allowed" (back-compat)."""
    import fitz

    doc = fitz.open(stream=body, filetype="pdf")
    try:
        perm = _perm_mask(allowed if allowed is not None else _ALL_PERMS)
        return doc.tobytes(
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=(owner_pw or user_pw),
            user_pw=user_pw,
            permissions=perm,
        )
    finally:
        doc.close()


def unlock_pdf(body: bytes, password: str) -> bytes:
    """Decrypt a PDF, given its correct password. Raises WrongPassword."""
    import fitz

    doc = fitz.open(stream=body, filetype="pdf")
    try:
        if doc.needs_pass and not doc.authenticate(password):
            raise WrongPassword()
        return doc.tobytes(encryption=fitz.PDF_ENCRYPT_NONE)
    finally:
        doc.close()


def compress(body: bytes, preset: str) -> bytes:
    """Compress a PDF with Ghostscript. Guarantees the result is never larger
    than the input — if compression doesn't help, the original is returned."""
    with tempfile.TemporaryDirectory(prefix="gs-", dir=TMP_ROOT) as work:
        in_path = os.path.join(work, "in.pdf")
        out_path = os.path.join(work, "out.pdf")
        with open(in_path, "wb") as fh:
            fh.write(body)
        ghostscript_compress(in_path, out_path, preset)
        with open(out_path, "rb") as fh:
            out = fh.read()
        return out if 0 < len(out) < len(body) else body


# ── MuPDF (PyMuPDF / fitz) in-place text editing ─────────────────────────────
# Extract real font/size/colour/bold per line; on save, redact ONLY the original
# text (leaving coloured backgrounds/graphics intact) and reinsert matched text.

_BASE14 = {
    ("sans", False, False): "helv", ("sans", True, False): "hebo",
    ("sans", False, True): "heit", ("sans", True, True): "hebi",
    ("serif", False, False): "tiro", ("serif", True, False): "tibo",
    ("serif", False, True): "tiit", ("serif", True, True): "tibi",
    ("mono", False, False): "cour", ("mono", True, False): "cobo",
    ("mono", False, True): "coit", ("mono", True, True): "cobi",
}


def _int_to_hex(srgb: int) -> str:
    return "#%06x" % (int(srgb) & 0xFFFFFF)


def _hex_to_rgb(h: str):
    h = (h or "#111111").lstrip("#")
    if len(h) != 6:
        return (0.07, 0.07, 0.07)
    return (int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255)


def _font_key(basefont: str) -> str:
    """A CSS-safe, stable id for an embedded font, normalised so the name a text
    span reports ("DejaVuSans-Bold") and the name get_page_fonts reports
    ("DejaVu Sans Bold") collapse to the SAME key. Strips the subset prefix
    (ABCDEF+) and all punctuation/spacing so the two always match — otherwise the
    saved text can't find its font and falls back to a base-14 substitute."""
    name = basefont or "font"
    if "+" in name:
        name = name.split("+", 1)[1]
    return "emb_" + (re.sub(r"[^A-Za-z0-9]", "", name).lower() or "font")


def _line_style(spans):
    s0 = spans[0]
    flags = int(s0.get("flags", 0))
    name = s0.get("font", "") or ""
    return {
        "text": "".join(s.get("text", "") for s in spans),
        "bbox": [
            min(s["bbox"][0] for s in spans), min(s["bbox"][1] for s in spans),
            max(s["bbox"][2] for s in spans), max(s["bbox"][3] for s in spans),
        ],
        "origin": [s0["origin"][0], s0["origin"][1]],
        "size": s0.get("size", 12),
        "color": _int_to_hex(s0.get("color", 0)),
        "bold": bool(flags & 16) or "bold" in name.lower() or "black" in name.lower(),
        "italic": bool(flags & 2) or "italic" in name.lower() or "oblique" in name.lower(),
        "family": _family_of(name, flags),
        # Embedded-font id (basefont) so the saved text reuses the exact font.
        "font": _font_key(name) if name else None,
    }


# Font names are a far more reliable signal than PyMuPDF's serif flag, which is
# often wrong (e.g. it flags DejaVu Sans as serif). Use the name first.
_SERIF_HINTS = ("times", "serif", "georgia", "roman", "garamond", "minion", "cambria",
                "book antiqua", "palatino", "baskerville", "caslon", "merriweather", "ming", "song")
_MONO_HINTS = ("mono", "courier", "consol", "menlo", "inconsolata", "source code")


def _family_of(name: str, flags: int) -> str:
    n = (name or "").lower()
    if any(k in n for k in _MONO_HINTS) or flags & 8:
        return "mono"
    if "sans" in n:  # an explicit sans name wins over a stray serif flag
        return "sans"
    if any(k in n for k in _SERIF_HINTS) or flags & 4:
        return "serif"
    return "sans"


def _extract_page_fonts(doc, page) -> dict:
    """Return {fontKey: {ext, b64}} for the page's EMBEDDED fonts so the browser
    can render the editor in the real fonts (and we can reuse them on save)."""
    out = {}
    try:
        page_fonts = doc.get_page_fonts(page.number)
    except Exception:
        return out
    for f in page_fonts:
        xref, basefont = f[0], f[3]
        if not basefont:
            continue
        key = _font_key(basefont)
        if key in out:
            continue
        try:
            _bf, ext, _ftype, buf = doc.extract_font(xref)
        except Exception:
            buf = b""
        if buf and ext in ("ttf", "otf", "cff", "woff", "woff2"):
            out[key] = {"ext": ext, "b64": base64.b64encode(buf).decode()}
    return out


def _redact_all_text(page) -> None:
    """Lift (remove) ALL text from a page, keeping graphics/images intact."""
    import fitz

    found = False
    for block in page.get_text("dict").get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            for s in line.get("spans", []):
                if s.get("text", "").strip():
                    page.add_redact_annot(fitz.Rect(s["bbox"]))
                    found = True
    if found:
        try:
            page.apply_redactions(
                images=fitz.PDF_REDACT_IMAGE_NONE, graphics=fitz.PDF_REDACT_LINE_ART_NONE
            )
        except TypeError:
            page.apply_redactions()


def _ocr_lines(page, dpi: int = 200) -> list:
    """OCR a scanned page → editable lines with boxes (Tesseract). Empty if OCR
    isn't available."""
    try:
        import io
        import pytesseract
        from PIL import Image
    except Exception:
        return []

    pix = page.get_pixmap(dpi=dpi)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    try:
        data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
    except Exception:
        return []

    scale = 72.0 / dpi
    groups: dict = {}
    for i in range(len(data["text"])):
        txt = (data["text"][i] or "").strip()
        try:
            conf = float(data["conf"][i])
        except (TypeError, ValueError):
            conf = -1
        if not txt or conf < 35:
            continue
        key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
        groups.setdefault(key, []).append((
            data["left"][i] * scale, data["top"][i] * scale,
            data["width"][i] * scale, data["height"][i] * scale, txt,
        ))

    lines = []
    for words in groups.values():
        x0 = min(w[0] for w in words)
        y0 = min(w[1] for w in words)
        x1 = max(w[0] + w[2] for w in words)
        y1 = max(w[1] + w[3] for w in words)
        h = y1 - y0
        lines.append({
            "text": " ".join(w[4] for w in words),
            "bbox": [x0, y0, x1, y1],
            "origin": [x0, y1 - h * 0.18],
            "size": max(6.0, h * 0.82),
            "color": "#111111", "bold": False, "italic": False, "family": "sans",
            "font": None, "ocr": True,
        })
    return lines


def _bg_hex(pix, bbox, scale) -> str:
    r, g, b = _sample_bg(pix, bbox, scale)
    return "#%02x%02x%02x" % (int(r * 255), int(g * 255), int(b * 255))


def edit_extract(pdf_bytes: bytes, page_index: int, dpi: int = 150) -> dict:
    """Return the page exactly as it renders (the original — so fonts are
    pixel-perfect) plus an invisible editable layer: one entry per text line with
    its real font/size/colour and the background colour behind it. Nothing is
    removed here — the original text stays visible until the user edits a line,
    and only edited lines are touched on save. Scanned pages are OCR'd the same way."""
    doc = __import__("fitz").open(stream=pdf_bytes, filetype="pdf")
    try:
        page = doc[page_index]
        rect = page.rect

        lines = []
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                spans = [s for s in line.get("spans", []) if s.get("text", "").strip()]
                if spans:
                    lines.append(_line_style(spans))

        scanned = not lines
        if scanned:  # no text layer → OCR the raster into editable lines
            lines = _ocr_lines(page)

        # The background is the ORIGINAL page (text included) — the editor only
        # covers a line when you actually edit it.
        pix = page.get_pixmap(dpi=dpi)
        sc = pix.width / rect.width if rect.width else 1
        for ln in lines:
            ln["bg"] = _bg_hex(pix, ln["bbox"], sc)

        return {
            "pageCount": doc.page_count, "width": rect.width, "height": rect.height,
            "image": base64.b64encode(pix.tobytes("png")).decode(),
            "lines": lines, "scanned": scanned,
        }
    finally:
        doc.close()


def _build_doc_fonts(doc) -> dict:
    """{fontKey: fitz.Font} for every embedded font, for exact reuse on save."""
    import fitz

    fonts = {}
    for pno in range(doc.page_count):
        try:
            page_fonts = doc.get_page_fonts(pno)
        except Exception:
            continue
        for f in page_fonts:
            key = _font_key(f[3])
            if not f[3] or key in fonts:
                continue
            try:
                _bf, _ext, _ft, buf = doc.extract_font(f[0])
                if buf:
                    fonts[key] = fitz.Font(fontbuffer=buf)
            except Exception:
                pass
    return fonts


def _draw_line(page, ln, doc_fonts, fallback_cache) -> None:
    """Draw one line with the EXACT embedded font where possible, falling back to
    a base-14 font per character for any glyph the (subset) font is missing."""
    import fitz

    text = ln.get("text") or ""
    if not text.strip():
        return
    size = float(ln.get("size", 12) or 12)
    color = _hex_to_rgb(ln.get("color", "#111111"))
    o = ln.get("origin") or [ln["bbox"][0], ln["bbox"][3] - size * 0.22]

    emb = doc_fonts.get(ln.get("font"))
    fb_key = (ln.get("family", "sans"), bool(ln.get("bold")), bool(ln.get("italic")))
    if fb_key not in fallback_cache:
        fallback_cache[fb_key] = fitz.Font(_BASE14.get(fb_key, "helv"))
    fb = fallback_cache[fb_key]

    # Group consecutive chars by which font can render them.
    runs = []
    for ch in text:
        if ch in (" ", "\t") and runs:
            f = runs[-1][0]
        elif emb is not None and emb.has_glyph(ord(ch)):
            f = emb
        else:
            f = fb
        if runs and runs[-1][0] is f:
            runs[-1][1] += ch
        else:
            runs.append([f, ch])

    tw = fitz.TextWriter(page.rect)
    pen = fitz.Point(o[0], o[1])
    try:
        for f, chunk in runs:
            try:
                pen = tw.append(pen, chunk, font=f, fontsize=size)
            except Exception:
                pen = tw.append(pen, chunk, font=fb, fontsize=size)
        tw.write_text(page, color=color)
    except Exception:
        page.insert_text(fitz.Point(o[0], o[1]), text, fontname=_BASE14.get(fb_key, "helv"),
                         fontsize=size, color=color)


def _pad(bbox, dx: float = 0.6, dy: float = 0.8):
    """Pad a line box slightly so redaction/cover catches glyph edges & descenders."""
    import fitz

    return fitz.Rect(bbox[0] - dx, bbox[1] - dy, bbox[2] + dx, bbox[3] + dy)


def _sample_bg(pix, bbox, scale):
    """Sample a background colour just left of a box (for covering scanned text)."""
    try:
        x = max(0, int(bbox[0] * scale) - 3)
        y = max(0, int(((bbox[1] + bbox[3]) / 2) * scale))
        r, g, b = pix.pixel(min(x, pix.width - 1), min(y, pix.height - 1))[:3]
        return (r / 255, g / 255, b / 255)
    except Exception:
        return (1, 1, 1)


def edit_apply(pdf_bytes: bytes, pages: list) -> bytes:
    """Apply ONLY the edited lines, leaving the rest of the page pristine.
    Digital pages: redact each edited line's box (removes the original glyphs with
    no ghost, keeps the background), then redraw it in the EXACT embedded font.
    Scanned pages: cover the box with the sampled background and draw the new text.
    Unedited lines are never sent, so unedited text stays original vector text."""
    import fitz

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        doc_fonts = _build_doc_fonts(doc)
        fallback_cache: dict = {}
        for pg in pages:
            idx = int(pg.get("page", 0))
            if idx < 0 or idx >= doc.page_count:
                continue
            page = doc[idx]
            edited = pg.get("lines", [])  # client sends only edited lines (empty text = erase)
            if not edited:
                continue
            scanned = not page.get_text("text").strip()

            if scanned:  # raster page — cover the word, then draw the new text
                pix = page.get_pixmap(dpi=150)
                sc = pix.width / page.rect.width if page.rect.width else 1
                for ln in edited:
                    fill = _sample_bg(pix, ln["bbox"], sc)
                    page.draw_rect(_pad(ln["bbox"]), color=fill, fill=fill)
                    _draw_line(page, ln, doc_fonts, fallback_cache)
            else:  # vector page — physically remove the original glyphs first
                for ln in edited:
                    page.add_redact_annot(_pad(ln["bbox"]))
                try:
                    page.apply_redactions(
                        images=fitz.PDF_REDACT_IMAGE_NONE, graphics=fitz.PDF_REDACT_LINE_ART_NONE
                    )
                except TypeError:
                    page.apply_redactions()
                for ln in edited:
                    _draw_line(page, ln, doc_fonts, fallback_cache)
        return doc.tobytes(deflate=True)
    finally:
        doc.close()


def redact_pdf(pdf_bytes: bytes, pages: list) -> bytes:
    """TRUE redaction: for each box, draw an opaque black rectangle AND remove the
    underlying text/vectors/image content beneath it (PyMuPDF apply_redactions) —
    so the hidden content can't be recovered by selecting or copying. Boxes are in
    PDF points with a top-left origin, matching the editor's coordinate space."""
    import fitz

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        for item in pages:
            idx = int(item.get("page", 0))
            if idx < 0 or idx >= doc.page_count:
                continue
            page = doc[idx]
            applied = False
            for b in item.get("boxes", []):
                try:
                    x, y, w, h = float(b["x"]), float(b["y"]), float(b["w"]), float(b["h"])
                except (KeyError, TypeError, ValueError):
                    continue
                if w <= 0 or h <= 0:
                    continue
                page.add_redact_annot(fitz.Rect(x, y, x + w, y + h), fill=(0, 0, 0))
                applied = True
            if applied:
                try:
                    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)
                except TypeError:
                    page.apply_redactions()
        return doc.tobytes(deflate=True, garbage=3)
    finally:
        doc.close()


def convert(body: bytes, source_ext: str, target: str, use_ai: bool = False) -> bytes:
    with tempfile.TemporaryDirectory(prefix="job-", dir=TMP_ROOT) as work:
        in_path = os.path.join(work, f"in.{source_ext}")
        with open(in_path, "wb") as fh:
            fh.write(body)
        out_path = os.path.join(work, f"out.{target}")

        if source_ext == "pdf" and target == "docx":
            pdf_to_docx(in_path, out_path)
        elif source_ext == "pdf" and target == "pptx":
            pdf_to_pptx(in_path, out_path, work, use_ai=use_ai)
        elif source_ext == "pdf" and target == "xlsx":
            pdf_to_xlsx(in_path, out_path)
        elif source_ext == "pdf" and target == "txt":
            pdf_to_txt(in_path, out_path)
        elif source_ext == "pdf" and target == "pdf":
            return body
        else:
            # Office → PDF (or Office → Office) via LibreOffice.
            produced = office_to_pdf_or_office(in_path, work, target)
            with open(produced, "rb") as fh:
                return fh.read()

        with open(out_path, "rb") as fh:
            return fh.read()


class Handler(BaseHTTPRequestHandler):
    def _send(self, code, body=b"", ctype="text/plain"):
        self.send_response(code)
        for k, v in CORS.items():
            self.send_header(k, v)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        if body:
            self.wfile.write(body)

    def do_OPTIONS(self):
        self._send(204)

    def do_GET(self):
        if self.path.startswith("/health"):
            self._send(200, b"ok")
        elif self.path.startswith("/ai-status"):
            active = fetch_active_model()
            if active:
                self._send(200, json.dumps({"available": True, "model": active["model"]}).encode(), "application/json")
            else:
                self._send(200, json.dumps({"available": ollama_available(), "model": OLLAMA_MODEL}).encode(), "application/json")
        else:
            self._send(404, b"Not found")

    def do_POST(self):
        parsed = urlparse(self.path)
        query = parse_qs(parsed.query)

        length = int(self.headers.get("Content-Length", "0"))
        if length > MAX_BODY_BYTES:
            return self._send(413, f"File too large (limit {MAX_BODY_BYTES // (1024 * 1024)} MB)".encode())
        body = self.rfile.read(length) if length else b""
        if not body:
            return self._send(400, b"Empty body")

        # URL-decoded password header (never the query string, so it can't leak
        # into access logs anywhere along the path).
        from urllib.parse import unquote
        password = unquote(self.headers.get("x-password", "") or "")

        try:
            # Bound concurrent heavy jobs so parallel uploads can't OOM the host.
            with _JOBS:
                if parsed.path == "/compress":
                    preset = (query.get("preset", ["ebook"])[0] or "ebook").lower()
                    out = compress(body, preset)
                    return self._send(200, out, "application/pdf")

                if parsed.path == "/protect":
                    if len(password) < 4:
                        return self._send(400, b"Password too short")
                    owner = unquote(self.headers.get("x-owner-password", "") or "")
                    # Absent header → all permissions (back-compat). Present →
                    # only the listed actions are allowed.
                    perms_header = self.headers.get("x-permissions", None)
                    allowed = (
                        {p.strip().lower() for p in perms_header.split(",") if p.strip()}
                        if perms_header is not None
                        else None
                    )
                    out = protect_pdf(body, password, owner_pw=(owner or None), allowed=allowed)
                    return self._send(200, out, "application/pdf")

                if parsed.path == "/unlock":
                    if not password:
                        return self._send(400, b"Password required")
                    out = unlock_pdf(body, password)
                    return self._send(200, out, "application/pdf")

                if parsed.path == "/convert":
                    target = (query.get("target", [""])[0] or "").lower()
                    source_ext = (self.headers.get("x-source-ext", "pdf") or "pdf").lower()
                    use_ai = query.get("ai", ["0"])[0] in ("1", "true", "yes")
                    if not target.isalnum() or len(target) > 5:
                        return self._send(400, b"Invalid target")
                    out = convert(body, source_ext, target, use_ai=use_ai)
                    return self._send(200, out, MIME.get(target, "application/octet-stream"))

                # In-place editing (MuPDF). /edit/page: raw PDF body + ?n=<index>.
                if parsed.path == "/edit/page":
                    n = int(query.get("n", ["0"])[0] or 0)
                    result = edit_extract(body, n)
                    return self._send(200, json.dumps(result).encode(), "application/json")

                # /edit/apply: JSON { pdf: base64, pages: [{page, lines:[...]}] } → edited PDF.
                if parsed.path == "/edit/apply":
                    payload = json.loads(body)
                    pdf = base64.b64decode(payload["pdf"])
                    out = edit_apply(pdf, payload.get("pages", []))
                    return self._send(200, out, "application/pdf")

                # /redact: JSON { pdf: base64, pages: [{page, boxes:[{x,y,w,h}]}] } → redacted PDF.
                if parsed.path == "/redact":
                    payload = json.loads(body)
                    pdf = base64.b64decode(payload["pdf"])
                    out = redact_pdf(pdf, payload.get("pages", []))
                    return self._send(200, out, "application/pdf")

            return self._send(404, b"Not found")
        except WrongPassword:
            self._send(401, b"Wrong password")
        except subprocess.CalledProcessError as exc:
            msg = exc.stderr.decode("utf-8", "replace") if exc.stderr else str(exc)
            self._send(500, f"Operation failed: {msg}".encode())
        except Exception as exc:  # noqa: BLE001
            self._send(500, f"Operation failed: {exc}".encode())

    def log_message(self, *args):
        pass


if __name__ == "__main__":
    print(f"PDFShell convert service on :{PORT} (tmp: {TMP_ROOT})", flush=True)
    threading.Thread(target=_janitor, daemon=True).start()
    ThreadingHTTPServer(("", PORT), Handler).serve_forever()
